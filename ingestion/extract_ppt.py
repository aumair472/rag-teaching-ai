"""
PowerPoint text extraction module.

Uses python-pptx to extract slide-level text with metadata.
"""

from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation

from app.core.logging import get_logger

logger = get_logger(__name__)


class PPTExtractor:
    """
    Extracts text from PowerPoint presentations with slide-level metadata.

    Iterates over all shapes on each slide and concatenates text frames.
    """

    def extract(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from a .pptx file.

        Args:
            file_path: Absolute or relative path to the PowerPoint file.

        Returns:
            A list of dictionaries, one per non-empty slide::

                {
                    "text": "...",
                    "slide": 1,
                    "source_name": "lecture_slides",
                    "source_type": "ppt"
                }

        Raises:
            FileNotFoundError: If the file does not exist.
        """
        path = Path(file_path).resolve()
        if not path.exists():
            raise FileNotFoundError(f"PPT not found: {path}")

        source_name = path.stem
        slides: List[Dict[str, Any]] = []

        logger.info(
            "Extracting PPT", extra={"source": source_name, "path": str(path)}
        )

        prs = Presentation(str(path))

        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts: List[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            text_parts.append(para_text)

                # Also extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                text_parts.append(cell_text)

            full_text = "\n".join(text_parts)

            if full_text.strip():
                slides.append({
                    "text": full_text,
                    "slide": slide_num,
                    "source_name": source_name,
                    "source_type": "ppt",
                })

        logger.info(
            "PPT extraction complete",
            extra={"source": source_name, "slides_extracted": len(slides)},
        )
        return slides
